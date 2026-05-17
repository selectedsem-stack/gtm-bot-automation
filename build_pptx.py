from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── Palette ────────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x1B, 0x2A, 0x4A)
BLUE_MID  = RGBColor(0x22, 0x3A, 0x5E)
ORANGE    = RGBColor(0xF4, 0xA6, 0x1D)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xD4, 0xDF, 0xF0)
GREEN     = RGBColor(0x27, 0xAE, 0x60)
GRAY_DARK = RGBColor(0x33, 0x33, 0x44)
GRAY_LIGHT= RGBColor(0xF0, 0xF4, 0xFA)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank


# ── Helpers ────────────────────────────────────────────────────────────────────

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape


def txbox(slide, l, t, w, h, text, size, bold=False, color=WHITE,
          align=PP_ALIGN.RIGHT, rtl=True, wrap=True, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    if rtl:
        _set_rtl(p)
    return box


def _set_rtl(para):
    pPr = para._pPr
    if pPr is None:
        pPr = para._p.get_or_add_pPr()
    pPr.set(qn("w:bidi"), "1")
    pPr.set("rtl", "1")


def bullet_slide(slide, items, l, t, w, size=20, color=WHITE,
                 bullet="•  ", gap=Pt(6), bold_first=False):
    box = slide.shapes.add_textbox(l, t, w, Inches(6))
    tf  = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.RIGHT
        p.space_after = gap
        run = p.add_run()
        run.text = bullet + item
        run.font.size  = Pt(size)
        run.font.color.rgb = color
        run.font.bold  = bold_first and (item == items[0])
        _set_rtl(p)
    return box


def slide_header(slide, title, subtitle=None):
    rect(slide, 0, 0, W, Inches(1.35), NAVY)
    rect(slide, 0, Inches(1.35), Inches(0.08), H - Inches(1.35), ORANGE)
    txbox(slide, Inches(0.3), Inches(0.18), W - Inches(0.6), Inches(0.75),
          title, 36, bold=True, color=WHITE)
    if subtitle:
        txbox(slide, Inches(0.3), Inches(0.92), W - Inches(0.6), Inches(0.45),
              subtitle, 18, color=ORANGE)


def tag_chip(slide, l, t, text, w=Inches(3.8), bg_color=BLUE_MID):
    r = rect(slide, l, t, w, Inches(0.48), bg_color)
    txbox(slide, l + Inches(0.1), t + Inches(0.04), w - Inches(0.2), Inches(0.44),
          text, 15, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)

# decorative orange bar bottom
rect(s, 0, H - Inches(0.18), W, Inches(0.18), ORANGE)
# decorative accent top-right corner block
rect(s, W - Inches(2.4), 0, Inches(2.4), Inches(0.1), ORANGE)

# logo placeholder circle
rect(s, Inches(0.55), Inches(0.55), Inches(1.4), Inches(1.4), BLUE_MID)
txbox(s, Inches(0.55), Inches(0.65), Inches(1.4), Inches(1.2),
      "SEL", 32, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, rtl=False)

txbox(s, Inches(0.3), Inches(2.2), W - Inches(0.6), Inches(1.3),
      "מערכת אוטומציה", 52, bold=True, color=WHITE)
txbox(s, Inches(0.3), Inches(3.3), W - Inches(0.6), Inches(1.0),
      "GTM + GA4", 52, bold=True, color=ORANGE)
txbox(s, Inches(0.3), Inches(4.5), W - Inches(0.6), Inches(0.6),
      "Selected Digital Agency", 22, color=LIGHT, italic=True)
txbox(s, Inches(0.3), Inches(5.1), W - Inches(0.6), Inches(0.5),
      "סיכום שלב א׳ | מאי 2026", 18, color=LIGHT)
txbox(s, Inches(0.3), Inches(5.75), W - Inches(0.6), Inches(0.5),
      "בניית האפליקציה — מהרעיון לייצור", 16, color=ORANGE, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — הבעיה
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, GRAY_LIGHT)
slide_header(s, "הבעיה שפתרנו", "כל לקוח חדש = שעות של עבודה ידנית")

# Before box
rect(s, Inches(0.4), Inches(1.7), Inches(5.6), Inches(5.0), RGBColor(0xFF,0xEB,0xEB))
txbox(s, Inches(0.4), Inches(1.75), Inches(5.6), Inches(0.5),
      "לפני — ידני", 20, bold=True, color=RGBColor(0xCC,0x00,0x00),
      align=PP_ALIGN.CENTER, rtl=False)
rect(s, Inches(0.4), Inches(2.2), Inches(5.6), Inches(0.04), RGBColor(0xCC,0x00,0x00))

before_items = [
    "כניסה ידנית ל-GTM — יצירת Container",
    "פתיחת GA4 — יצירת Property ידנית",
    "הוספת תגים אחד-אחד (10+ תגים)",
    "טריגרים ידניים לכל אירוע",
    "העתקת קודי script ל-WordPress",
    "בדיקה, תיקונים, בדיקה שוב...",
    "זמן: 2–4 שעות לכל לקוח",
]
bullet_slide(s, before_items, Inches(0.5), Inches(2.35), Inches(5.3),
             size=17, color=GRAY_DARK, bullet="✗  ")

# After box
rect(s, Inches(7.1), Inches(1.7), Inches(5.6), Inches(5.0), RGBColor(0xE8,0xF8,0xF0))
txbox(s, Inches(7.1), Inches(1.75), Inches(5.6), Inches(0.5),
      "אחרי — אוטומטי", 20, bold=True, color=GREEN,
      align=PP_ALIGN.CENTER, rtl=False)
rect(s, Inches(7.1), Inches(2.2), Inches(5.6), Inches(0.04), GREEN)

after_items = [
    "ממשק Web פשוט — 3 שלבים",
    "GTM Container — נוצר אוטומטי",
    "GA4 Property + Stream — נוצרים אוטומטי",
    "כל התגים — מיובאים אוטומטי",
    "קוד HTML — מוכן להדבקה",
    "תיעוד מלא של כל ה-IDs",
    "זמן: 30–40 שניות",
]
bullet_slide(s, after_items, Inches(7.2), Inches(2.35), Inches(5.3),
             size=17, color=GRAY_DARK, bullet="✓  ")

# Arrow in middle
txbox(s, Inches(5.8), Inches(3.7), Inches(1.7), Inches(0.8),
      "→", 60, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, rtl=False)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — מה בנינו
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, GRAY_LIGHT)
slide_header(s, "מה בנינו", "אפליקציית Web פנימית לאוטומציה מלאה")

tech_items = [
    ("Python + Streamlit", "שפה וממשק ה-Web", NAVY),
    ("GitHub", "ניהול קוד ופריסה אוטומטית", RGBColor(0x24, 0x29, 0x2F)),
    ("Streamlit Cloud", "אחסון ופריסה בענן — HTTPS", RGBColor(0xFF, 0x4B, 0x4B)),
    ("GTM API v2", "יצירת containers ותגים", RGBColor(0x34, 0x65, 0xA4)),
    ("GA4 Analytics Admin API", "יצירת Properties ו-Streams", RGBColor(0xE3, 0x74, 0x00)),
    ("OAuth 2.0", "אימות מאובטח מול גוגל", RGBColor(0x0F, 0x9D, 0x58)),
]

cols = [(Inches(0.4), Inches(0.45)), (Inches(4.6), Inches(0.45)), (Inches(8.8), Inches(0.45))]
row_h = Inches(1.5)
for i, (title, desc, color) in enumerate(tech_items):
    col_idx = i % 3
    row_idx = i // 3
    lx = cols[col_idx][0]
    ty = Inches(1.7) + row_idx * row_h + cols[col_idx][1]
    w  = Inches(3.8)
    h  = Inches(1.3)
    r = rect(s, lx, ty, w, h, color)
    txbox(s, lx + Inches(0.1), ty + Inches(0.08), w - Inches(0.2), Inches(0.5),
          title, 18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, rtl=False)
    txbox(s, lx + Inches(0.1), ty + Inches(0.58), w - Inches(0.2), Inches(0.6),
          desc, 14, color=LIGHT, align=PP_ALIGN.CENTER, rtl=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — APIs של גוגל
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, GRAY_LIGHT)
slide_header(s, "ההתממשקות ל-API של גוגל", "שני שירותי Google פועלים במקביל")

# GTM card
rect(s, Inches(0.4), Inches(1.65), Inches(6.0), Inches(5.25), NAVY)
txbox(s, Inches(0.5), Inches(1.75), Inches(5.8), Inches(0.55),
      "GTM API v2", 24, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, rtl=False)
rect(s, Inches(0.4), Inches(2.25), Inches(6.0), Inches(0.04), ORANGE)
gtm_actions = [
    "יצירת Container חדש",
    "קריאת workspace",
    "יצירת Built-in Variables",
    "יצירת Triggers",
    "יצירת Variables",
    "יצירת Tags",
    "יצירת גרסה + פרסום",
]
bullet_slide(s, gtm_actions, Inches(0.5), Inches(2.35), Inches(5.7),
             size=17, color=LIGHT, bullet="▸  ")

# GA4 card
rect(s, Inches(7.0), Inches(1.65), Inches(6.0), Inches(5.25), NAVY)
txbox(s, Inches(7.1), Inches(1.75), Inches(5.8), Inches(0.55),
      "Analytics Admin API", 24, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, rtl=False)
rect(s, Inches(7.0), Inches(2.25), Inches(6.0), Inches(0.04), ORANGE)
ga4_actions = [
    "שליפת כל ה-GA4 Accounts",
    "יצירת GA4 Property",
    "הגדרות: שפה, מטבע, timezone",
    "יצירת Web Data Stream",
    "שליפת Measurement ID (G-XXXX)",
    "שליפת Stream ID",
    "שליפת Property ID",
]
bullet_slide(s, ga4_actions, Inches(7.1), Inches(2.35), Inches(5.7),
             size=17, color=LIGHT, bullet="▸  ")

# OAuth note
rect(s, Inches(0.4), H - Inches(0.65), W - Inches(0.8), Inches(0.5),
     RGBColor(0xE8, 0xF4, 0xFF))
txbox(s, Inches(0.5), H - Inches(0.65), W - Inches(1.0), Inches(0.5),
      "OAuth 2.0 — Token מוחזק ב-Streamlit Secrets | מתרענן אוטומטית בכל טעינה",
      14, color=NAVY, align=PP_ALIGN.CENTER, rtl=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — 4 שלבי ההקמה
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, GRAY_LIGHT)
slide_header(s, "תהליך ההקמה — 4 שלבים", "חוויית משתמש פשוטה, לוגיקה מורכבת מאחורי הקלעים")

steps = [
    ("1", "פרטי הלקוח", ["שם לקוח + דומיין", "עמוד תודה לטופס", "האם חנות eCommerce?", "Google Ads / Facebook Pixel / Maskyoo"]),
    ("2", "בחירת GA4 Account", ["האפליקציה שולפת את כל 15 חשבונות GA4", "המשתמש בוחר לאיזה חשבון משייכים", "מגבלה: עד 100 Properties לחשבון"]),
    ("3", "אישור ובדיקה", ["סיכום כל הפרטים לפני הרצה", "בדיקה אחרונה לפני הפעלה", "חזרה לתיקון אם צריך"]),
    ("4", "הרצה ותוצאות", ["30–40 שניות של אוטומציה", "לוג חי על המסך", "תצוגת IDs + תגים + קודי HTML"]),
]

step_colors = [NAVY, BLUE_MID, RGBColor(0x1A,0x4A,0x7A), RGBColor(0x0D,0x3B,0x66)]
sw = Inches(3.0)
for i, (num, title, bullets) in enumerate(steps):
    lx = Inches(0.25) + i * Inches(3.25)
    ty = Inches(1.65)
    rect(s, lx, ty, sw, Inches(5.3), step_colors[i])
    # number circle
    rect(s, lx + Inches(1.0), ty + Inches(0.15), Inches(1.0), Inches(0.85), ORANGE)
    txbox(s, lx + Inches(1.0), ty + Inches(0.15), Inches(1.0), Inches(0.85),
          num, 36, bold=True, color=WHITE, align=PP_ALIGN.CENTER, rtl=False)
    txbox(s, lx + Inches(0.1), ty + Inches(1.1), sw - Inches(0.2), Inches(0.55),
          title, 19, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    rect(s, lx + Inches(0.1), ty + Inches(1.62), sw - Inches(0.2), Inches(0.04), ORANGE)
    bullet_slide(s, bullets, lx + Inches(0.1), ty + Inches(1.75), sw - Inches(0.2),
                 size=15, color=LIGHT, bullet="· ")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — מה קורה מאחורי הקלעים
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
slide_header(s, "מה קורה מאחורי הקלעים", "9 פעולות אוטומטיות בסדר מדויק")

actions = [
    ("1", "יצירת GTM Container", "שם = דומיין הלקוח | Platform: Web"),
    ("2", "יצירת GA4 Property", "שם לקוח | Timezone: Jerusalem | מטבע: ILS"),
    ("3", "יצירת Web Data Stream", "URI: https://דומיין"),
    ("4", "קריאת template.json", "תבנית התגים הפנימית של הסוכנות"),
    ("5", "החלפת PLACEHOLDERs", "GA4 ID, Stream ID, Domain, Thank-you URL, Ads/Pixel/Maskyoo"),
    ("6", "סינון תגים", "הסרת תגים לא-רלוונטיים (eCommerce / Ads / Pixel)"),
    ("7", "ייבוא Triggers + Variables", "כל הטריגרים והמשתנים נוצרים בסדר נכון"),
    ("8", "ייבוא Tags", "כל התגים עם ה-trigger IDs החדשים"),
    ("9", "פרסום גרסה", "\"Initial setup — שם לקוח — תאריך\""),
]

col1_items = actions[:5]
col2_items = actions[5:]

for i, (num, title, desc) in enumerate(col1_items):
    ty = Inches(1.7) + i * Inches(1.0)
    rect(s, Inches(0.3), ty, Inches(0.55), Inches(0.55), ORANGE)
    txbox(s, Inches(0.3), ty, Inches(0.55), Inches(0.55),
          num, 20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, rtl=False)
    txbox(s, Inches(1.0), ty, Inches(5.5), Inches(0.35),
          title, 17, bold=True, color=WHITE)
    txbox(s, Inches(1.0), ty + Inches(0.35), Inches(5.5), Inches(0.35),
          desc, 14, color=LIGHT)

for i, (num, title, desc) in enumerate(col2_items):
    ty = Inches(1.7) + i * Inches(1.0)
    rect(s, Inches(6.8), ty, Inches(0.55), Inches(0.55), ORANGE)
    txbox(s, Inches(6.8), ty, Inches(0.55), Inches(0.55),
          num, 20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, rtl=False)
    txbox(s, Inches(7.5), ty, Inches(5.5), Inches(0.35),
          title, 17, bold=True, color=WHITE)
    txbox(s, Inches(7.5), ty + Inches(0.35), Inches(5.5), Inches(0.35),
          desc, 14, color=LIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — התגים הקבועים
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, GRAY_LIGHT)
slide_header(s, "התגים שמוקמים — תמיד", "7 תגים קבועים בכל לקוח, ללא יוצא מן הכלל")

always_tags = [
    ("GA4 — Configuration", "Google Tag", "Measurement ID הראשי | כל עמוד"),
    ("GA4 — Stream ID", "Google Tag", "Stream ID נפרד | כל עמוד"),
    ("Conversion Linker", "Conversion Linker", "קישור המרות Ads ↔ Analytics | כל עמוד"),
    ("GA4 — Event — page_view", "GA4 Event", "מדידת צפיות בעמוד | כל עמוד"),
    ("GA4 — Event — contact_form_submission", "GA4 Event", "טופס יצירת קשר | הגעה לעמוד תודה"),
    ("GA4 — Event — phone_clicks", "GA4 Event", "לחיצה על מספר טלפון"),
    ("GA4 — Event — wa_clicks", "GA4 Event", "לחיצה על WhatsApp"),
]

tag_colors = [
    RGBColor(0x23,0x47,0x8A), RGBColor(0x1D,0x60,0xAA), RGBColor(0x15,0x75,0xC2),
    RGBColor(0x0D,0x86,0xD4), RGBColor(0x05,0x93,0xDD), RGBColor(0x00,0x9E,0xE3),
    RGBColor(0x00,0xA8,0xE8),
]

cols_t = [Inches(0.3), Inches(4.55), Inches(8.8)]
rows_t = [Inches(1.7), Inches(2.7), Inches(3.7), Inches(4.7)]

for i, (name, tag_type, desc) in enumerate(always_tags):
    col = i % 3 if i < 6 else 1
    row = i // 3
    if i == 6:
        lx = Inches(4.55)
        ty = Inches(4.7)
    else:
        lx = cols_t[i % 3]
        ty = rows_t[i // 3]
    tw = Inches(3.8)
    rect(s, lx, ty, tw, Inches(0.85), tag_colors[i])
    txbox(s, lx + Inches(0.1), ty + Inches(0.05), tw - Inches(0.2), Inches(0.38),
          name, 16, bold=True, color=WHITE)
    txbox(s, lx + Inches(0.1), ty + Inches(0.44), tw - Inches(0.2), Inches(0.35),
          desc, 12, color=LIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — תגים אופציונליים
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, GRAY_LIGHT)
slide_header(s, "תגים אופציונליים — לפי הלקוח", "נוספים רק אם נבחרו בטופס, אחרת מסוננים")

groups = [
    (
        "חנות eCommerce", RGBColor(0x6A,0x0D,0xAD),
        ["GA4 — Event — purchase", "GA4 — Event — add_to_cart",
         "GA4 — Event — remove_from_cart", "GA4 — Event — begin_checkout",
         "GA4 — Event — view_item"],
        "+5 תגים | triggers ל-WooCommerce dataLayer"
    ),
    (
        "Google Ads", RGBColor(0x1A,0x73,0xE8),
        ["Google Ads — Remarketing"],
        "+1 תג | Conversion ID: AW-XXXXXXXXX"
    ),
    (
        "Facebook Pixel", RGBColor(0x18,0x77,0xF2),
        ["Facebook Pixel — Base (PageView)", "Facebook Pixel — Event — purchase"],
        "+2 תגים | Pixel ID: 15–16 ספרות"
    ),
    (
        "Maskyoo", RGBColor(0x00,0x78,0x56),
        ["Maskyoo Script — Virtual Number"],
        "+1 תג | מספר וירטואלי + חיבור ל-GA4 ID"
    ),
]

positions = [
    (Inches(0.3), Inches(1.65), Inches(6.15)),
    (Inches(6.9), Inches(1.65), Inches(6.1)),
    (Inches(0.3), Inches(4.25), Inches(6.15)),
    (Inches(6.9), Inches(4.25), Inches(6.1)),
]

for (title, color, tags, note), (lx, ty, gw) in zip(groups, positions):
    gh = Inches(2.4)
    rect(s, lx, ty, gw, gh, color)
    txbox(s, lx + Inches(0.15), ty + Inches(0.1), gw - Inches(0.3), Inches(0.45),
          title, 20, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    rect(s, lx, ty + Inches(0.55), gw, Inches(0.04), WHITE)
    bullet_slide(s, tags, lx + Inches(0.15), ty + Inches(0.65), gw - Inches(0.3),
                 size=15, color=WHITE, bullet="✅  ")
    rect(s, lx, ty + gh - Inches(0.45), gw, Inches(0.45), RGBColor(0,0,0))
    txbox(s, lx + Inches(0.1), ty + gh - Inches(0.44), gw - Inches(0.2), Inches(0.42),
          note, 13, color=ORANGE, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — template.json
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
slide_header(s, "template.json — לב המערכת", "המדריך הידני של הסוכנות, מקודד פעם אחת לתבנית JSON")

txbox(s, Inches(0.4), Inches(1.6), Inches(7.5), Inches(0.55),
      "מה זה template.json?", 22, bold=True, color=ORANGE)
txbox(s, Inches(0.4), Inches(2.15), Inches(7.5), Inches(1.5),
      "כל תג, טריגר ומשתנה שהיינו מגדירים ידנית בכל לקוח — מקודד אחת ולתמיד בקובץ JSON אחד.\n"
      "האפליקציה קוראת אותו, מחליפה PLACEHOLDERs ומייבאת ישירות ל-GTM.",
      16, color=LIGHT)

placeholders = [
    ("PLACEHOLDER_GA4_ID", "GA4 Measurement ID", "G-XXXXXXXXXX"),
    ("PLACEHOLDER_STREAM_ID", "GA4 Stream ID", "מספר ארוך"),
    ("PLACEHOLDER_ADS_ID", "Google Ads Conversion ID", "AW-XXXXXXXXX"),
    ("PLACEHOLDER_PIXEL_ID", "Facebook Pixel ID", "15–16 ספרות"),
    ("PLACEHOLDER_DOMAIN", "דומיין האתר", "example.co.il"),
    ("PLACEHOLDER_THANKYOU_URL", "עמוד תודה", "/thank-you/"),
    ("PLACEHOLDER_MASKYOO_NUMBER", "מספר Maskyoo", "055-XXXXXXX"),
]

ty_start = Inches(3.75)
for i, (ph, name, fmt) in enumerate(placeholders):
    ty = ty_start + i * Inches(0.5)
    rect(s, Inches(0.4), ty, Inches(4.2), Inches(0.42), RGBColor(0x12,0x20,0x40))
    txbox(s, Inches(0.5), ty + Inches(0.03), Inches(4.0), Inches(0.38),
          ph, 13, color=ORANGE, rtl=False, align=PP_ALIGN.LEFT)
    txbox(s, Inches(4.8), ty + Inches(0.03), Inches(3.5), Inches(0.38),
          name, 14, color=WHITE)
    txbox(s, Inches(10.2), ty + Inches(0.03), Inches(2.8), Inches(0.38),
          fmt, 13, color=LIGHT, italic=True)

txbox(s, Inches(0.4), Inches(1.6), Inches(12.5), Inches(0.45),
      "PLACEHOLDERs — 7 ערכים משתנים", 14, bold=True,
      color=GRAY_LIGHT, align=PP_ALIGN.LEFT, rtl=False)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — התוצר הסופי
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, GRAY_LIGHT)
slide_header(s, "התוצר הסופי", "מה רואים בסיום ההרצה — 30 שניות")

# 3 ID boxes
id_boxes = [
    ("GTM-XXXXXXX", "GTM Container ID", NAVY),
    ("G-XXXXXXXXXX", "GA4 Measurement ID", RGBColor(0xE3,0x74,0x00)),
    ("123456789", "GA4 Property ID", GREEN),
]
for i, (val, label, color) in enumerate(id_boxes):
    lx = Inches(0.4) + i * Inches(4.3)
    rect(s, lx, Inches(1.65), Inches(4.0), Inches(1.2), color)
    txbox(s, lx + Inches(0.1), Inches(1.72), Inches(3.8), Inches(0.55),
          val, 26, bold=True, color=WHITE, align=PP_ALIGN.CENTER, rtl=False)
    txbox(s, lx + Inches(0.1), Inches(2.2), Inches(3.8), Inches(0.4),
          label, 14, color=LIGHT, align=PP_ALIGN.CENTER)

# Deliverables
deliverables = [
    "רשימת כל התגים שהוקמו עם ✅ לכל תג",
    "קוד HEAD מוכן — הדבק מעל </head>",
    "קוד BODY מוכן — הדבק אחרי <body>",
    "לוג הרצה מלא לבדיקה",
    "קישור ישיר ל-CRM עם Property ID",
]
txbox(s, Inches(0.4), Inches(3.05), Inches(8.0), Inches(0.45),
      "מה מקבלים בנוסף:", 18, bold=True, color=NAVY)
bullet_slide(s, deliverables, Inches(0.4), Inches(3.5), Inches(8.0),
             size=18, color=GRAY_DARK, bullet="✅  ")

# Time badge
rect(s, Inches(9.5), Inches(3.1), Inches(3.5), Inches(3.5), NAVY)
txbox(s, Inches(9.5), Inches(3.3), Inches(3.5), Inches(0.9),
      "30–40", 54, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, rtl=False)
txbox(s, Inches(9.5), Inches(4.2), Inches(3.5), Inches(0.55),
      "שניות", 28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(s, Inches(9.5), Inches(4.75), Inches(3.5), Inches(0.45),
      "לקוח מוכן לאוויר", 16, color=LIGHT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — שלב הבא
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
slide_header(s, "שלב הבא — אוטומציה מלאה על האתרים", "סגירת המעגל: שתילת הסקריפט אוטומטית על WordPress")

# Current state
rect(s, Inches(0.3), Inches(1.65), Inches(5.8), Inches(5.3), BLUE_MID)
txbox(s, Inches(0.3), Inches(1.7), Inches(5.8), Inches(0.5),
      "שלב א׳ — הושלם ✅", 20, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
rect(s, Inches(0.3), Inches(2.17), Inches(5.8), Inches(0.04), GREEN)
done_items = [
    "GTM Container — אוטומטי",
    "GA4 Property + Stream — אוטומטי",
    "כל התגים — אוטומטי",
    "גרסה מפורסמת — אוטומטי",
    "קודי HTML מוכנים",
    "",
    "נשאר ידני: הדבקת HTML ב-WordPress",
]
bullet_slide(s, done_items, Inches(0.4), Inches(2.25), Inches(5.6),
             size=16, color=LIGHT, bullet="✅  ")

# Arrow
txbox(s, Inches(6.1), Inches(3.9), Inches(1.2), Inches(0.8),
      "→", 54, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, rtl=False)

# Next phase
rect(s, Inches(7.2), Inches(1.65), Inches(5.8), Inches(5.3), RGBColor(0x0D,0x3B,0x66))
txbox(s, Inches(7.2), Inches(1.7), Inches(5.8), Inches(0.5),
      "שלב ב׳ — בתכנון", 20, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
rect(s, Inches(7.2), Inches(2.17), Inches(5.8), Inches(0.04), ORANGE)
next_items = [
    "חיבור SSH לשרתי Cloudways",
    "זיהוי WordPress של הלקוח ב-sites.json",
    "WP-CLI: הוספת קוד HEAD + BODY",
    "אימות שהסקריפט אכן פועל",
    "לחלופין: REST API + Application Password",
    "לחלופין: פלאגין Insert Headers & Footers",
]
bullet_slide(s, next_items, Inches(7.3), Inches(2.25), Inches(5.6),
             size=16, color=LIGHT, bullet="→  ")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — תשתית שלב ב׳
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, GRAY_LIGHT)
slide_header(s, "תשתית שלב ב׳ — כבר קיימת", "כל הבסיס מוכן, נדרש רק לחבר אותו לאפליקציה")

infra_items = [
    ("sites.json", "קובץ מרכזי עם כל האתרים", ["כתובת שרת Cloudways", "נתיב WordPress", "Application Password", "פרטי SSH"], NAVY),
    ("SSH Config", "חיבור מוגדר לכל שרת", ["כינויים: cloudways-server1/2...", "מפתח SSH מוגדר", "גישה ישירה בלי סיסמה"], RGBColor(0x1A,0x4A,0x7A)),
    ("WP-CLI", "כלי שורת פקודה ל-WordPress", ["מותקן על כל שרת", "פקודות: option update, eval-file", "מאפשר שינויים ברמת DB"], RGBColor(0x21,0x59,0x32)),
    ("App Password", "גישת API ל-WordPress", ["מוגדר לכל אתר", "REST API /wp-json/", "לחלופין לגישה ללא SSH"], RGBColor(0x6A,0x1A,0x0D)),
]

for i, (title, subtitle, items, color) in enumerate(infra_items):
    lx = Inches(0.3) + i * Inches(3.25)
    ty = Inches(1.65)
    w  = Inches(3.05)
    rect(s, lx, ty, w, Inches(5.3), color)
    txbox(s, lx + Inches(0.1), ty + Inches(0.1), w - Inches(0.2), Inches(0.48),
          title, 20, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    rect(s, lx + Inches(0.1), ty + Inches(0.58), w - Inches(0.2), Inches(0.04), ORANGE)
    txbox(s, lx + Inches(0.1), ty + Inches(0.65), w - Inches(0.2), Inches(0.45),
          subtitle, 14, color=LIGHT, align=PP_ALIGN.CENTER, italic=True)
    bullet_slide(s, items, lx + Inches(0.1), ty + Inches(1.15), w - Inches(0.2),
                 size=14, color=WHITE, bullet="·  ")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — לפני / אחרי
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, GRAY_LIGHT)
slide_header(s, "לפני vs אחרי — הפוך לשחור-לבן", "מה השתנה עבור הצוות")

comparisons = [
    ("זמן הקמת לקוח", "2–4 שעות", "30–40 שניות"),
    ("תלות בידע אישי", "רק מי שיודע GTM", "כל אחד בצוות"),
    ("עקביות בתגים", "תלויה באדם", "זהה בכל פעם"),
    ("תיעוד", "ידני / לא קיים", "אוטומטי — IDs + תגים"),
    ("שגיאות אנוש", "אפשריות", "אפס"),
    ("הטמעה ב-WordPress", "ידנית", "קוד מוכן — עתידי: אוטומטי"),
    ("פרסום גרסה GTM", "ידני", "אוטומטי"),
    ("GA4 Property", "פתיחה ידנית", "אוטומטי עם כל הגדרות"),
]

# Headers
rect(s, Inches(0.3), Inches(1.6), Inches(3.8), Inches(0.5), NAVY)
txbox(s, Inches(0.3), Inches(1.6), Inches(3.8), Inches(0.5),
      "נושא", 16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rect(s, Inches(4.3), Inches(1.6), Inches(4.1), Inches(0.5), RGBColor(0xCC,0x00,0x00))
txbox(s, Inches(4.3), Inches(1.6), Inches(4.1), Inches(0.5),
      "לפני", 16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, rtl=False)
rect(s, Inches(8.6), Inches(1.6), Inches(4.5), Inches(0.5), GREEN)
txbox(s, Inches(8.6), Inches(1.6), Inches(4.5), Inches(0.5),
      "אחרי", 16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, rtl=False)

for i, (topic, before, after) in enumerate(comparisons):
    ty = Inches(2.12) + i * Inches(0.62)
    bg_c = GRAY_LIGHT if i % 2 == 0 else WHITE
    rect(s, Inches(0.3), ty, Inches(3.8), Inches(0.6), bg_c)
    rect(s, Inches(4.3), ty, Inches(4.1), Inches(0.6), RGBColor(0xFF,0xF0,0xF0))
    rect(s, Inches(8.6), ty, Inches(4.5), Inches(0.6), RGBColor(0xF0,0xFF,0xF4))
    txbox(s, Inches(0.4), ty + Inches(0.08), Inches(3.6), Inches(0.45),
          topic, 15, bold=True, color=NAVY)
    txbox(s, Inches(4.4), ty + Inches(0.08), Inches(3.9), Inches(0.45),
          before, 15, color=RGBColor(0xAA,0x00,0x00))
    txbox(s, Inches(8.7), ty + Inches(0.08), Inches(4.3), Inches(0.45),
          after, 15, color=RGBColor(0x14,0x6B,0x3A))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — סיכום + Next Steps
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
rect(s, 0, H - Inches(0.2), W, Inches(0.2), ORANGE)

txbox(s, Inches(0.5), Inches(0.4), W - Inches(1.0), Inches(0.65),
      "שלב א׳ הושלם בהצלחה", 38, bold=True, color=WHITE)
txbox(s, Inches(0.5), Inches(1.05), W - Inches(1.0), Inches(0.45),
      "מבקשים אישור להתקדם לשלב ב׳", 22, color=ORANGE, italic=True)

rect(s, Inches(0.3), Inches(1.7), Inches(5.9), Inches(4.8), BLUE_MID)
txbox(s, Inches(0.4), Inches(1.75), Inches(5.7), Inches(0.5),
      "מה הושג — שלב א׳", 20, bold=True, color=ORANGE)
rect(s, Inches(0.3), Inches(2.22), Inches(5.9), Inches(0.04), ORANGE)
achieved = [
    "אפליקציית Web פנימית — פרוסה ופעילה",
    "אוטומציה מלאה של GTM + GA4",
    "כל התגים — מוכנים ועקביים",
    "חיבור מאובטח ל-Google APIs",
    "ממשק פשוט — כל הצוות יכול להשתמש",
    "הפחתה מ-4 שעות → 40 שניות",
]
bullet_slide(s, achieved, Inches(0.4), Inches(2.3), Inches(5.7),
             size=17, color=WHITE, bullet="✅  ")

rect(s, Inches(7.1), Inches(1.7), Inches(5.9), Inches(4.8), RGBColor(0x0D,0x3B,0x66))
txbox(s, Inches(7.2), Inches(1.75), Inches(5.7), Inches(0.5),
      "הבקשה — שלב ב׳", 20, bold=True, color=ORANGE)
rect(s, Inches(7.1), Inches(2.22), Inches(5.9), Inches(0.04), ORANGE)
next_steps = [
    "אישור תקציב + זמן פיתוח לשלב ב׳",
    "חיבור SSH לשרתי Cloudways",
    "שתילת סקריפט GTM אוטומטי ב-WordPress",
    "בדיקה על אתר אחד → rollout לכל האתרים",
    "תוצאה: לקוח חדש חי — ללא מגע אנושי",
]
bullet_slide(s, next_steps, Inches(7.2), Inches(2.3), Inches(5.7),
             size=17, color=WHITE, bullet="→  ")

# Bottom tagline
txbox(s, Inches(0.5), H - Inches(0.75), W - Inches(1.0), Inches(0.5),
      "Selected Digital Agency | GTM + GA4 Automation | 2026",
      14, color=LIGHT, align=PP_ALIGN.CENTER, italic=True)

# ── Save ───────────────────────────────────────────────────────────────────────
out = r"C:\Users\SEMSELECTED\projects\gtm-bot-automation\GTM-GA4-Automation-Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
